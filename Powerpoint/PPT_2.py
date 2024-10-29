from pptx import Presentation
import pandas as pd

# Load the PowerPoint presentation
presentation = Presentation(r'/////////////////////////////////////')

# Identify shape combinations
combination_dict = {}
data_list=[]

for slide_number, slide in enumerate(presentation.slides, start=1):
    slide_data = {}  
    for shape_number, shape in enumerate(slide.shapes, start=1):
        if shape.has_text_frame:
            text=[]
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
            
            data=[]
            if shape.name=="TextBox 17":
                date_text="".join(text)
                slide_data['Date'] = date_text
                print(f"Date:{date_text}")
            if shape.name=="TextBox 5":
                title="".join(text)
                slide_data['Title'] = title
                print(f"Title: {title}")
                #print(f"Title: {text}")
            if shape.name=="TextBox 190":
                suppourt="".join(text)
                slide_data['Suppourt']=suppourt
                print(f"Suppourt Needed: {suppourt}")
            if shape.name=="TextBox 371":
                week_went_by="".join(text)
                slide_data["Week Went By"]=week_went_by
                print(f"Week Went By: {week_went_by}")
               #print(f"Week Went By: {text}")
            if shape.name=="TextBox 14":
                week_ahead="".join(text)
                slide_data["Week Ahead"]=tmss_week_ahead
                print(f"Week Ahead:{week_ahead}")
                #print(f"Week Ahead:{text}")
            if shape.name=="TextBox 370":
                week_went_by=''.join(text)
                slide_data["Week Went By"]=week_went_by
                print(f"Week Went By: {week_went_by}/n")
                #print(f"Week Went By: {text}")
            if shape.name=="TextBox 372":
                week_ahead="".join(text)
                slide_data[Week Ahead"]=sust_week_ahead
                print(f"Week Ahead: {week_ahead}")
            if shape.name=="TextBox 23":
                highlights="".join(text)
                slide_data["Hihlights"]=highlights
                print(f"Highlights:{highlights}") 
                #print(f"Highlights:{text}") 
           
            
          
                
    if slide_data:
        slide_data['Slide Number'] = slide_number  # Optionally add slide number
        data_list.append(slide_data)   

desired_order = ['Date', 'Title', 'Support',  'Highlights','Shipment Forecast', 'Slide Number']  #missing entries here
df = pd.DataFrame(data_list).T
df = df.reindex(columns=[col for col in desired_order if col in df.columns])

print(df)

df.to_csv("Report1.csv", index=False,encoding='utf-8',sep=",")
