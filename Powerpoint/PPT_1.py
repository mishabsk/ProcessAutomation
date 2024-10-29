from pptx import Presentation

def extract_text_from_slide(slide):
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text.append(run.text)
    return '\n'.join(text)

def extract_shapes_from_slide(slide):
    shapes_info = []
    for shape in slide.shapes:
        shape_info = {
            'shape_type': shape.shape_type,
            'text': shape.text if shape.has_text_frame else 'No text',
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height
        }
        shapes_info.append(shape_info)
    return shapes_info

def extract_headers_from_slide(slide):
    headers = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text_frame.text.startswith('Header'):
                headers.append(shape.text_frame.text)
    return headers

def extract_presentation_info(pptx_file):
    presentation = Presentation(pptx_file)
    all_texts = []
    all_shapes = []
    all_headers = []

    for slide_number, slide in enumerate(presentation.slides):
        print(f"Processing Slide {slide_number + 1}")

        text = extract_text_from_slide(slide)
        shapes = extract_shapes_from_slide(slide)
        headers = extract_headers_from_slide(slide)

        all_texts.append(text)
        all_shapes.append(shapes)
        all_headers.extend(headers)

    return all_texts, all_shapes, all_headers

if __name__ == "__main__":
    pptx_file = r"C///////////////////////"  # Path to your PowerPoint file
    texts, shapes, headers = extract_presentation_info(pptx_file)

    #print("Extracted Texts:")
    #for text in texts:
    #    print(text)
    
    print("\nExtracted Shapes:")
    for shape_list in shapes:
        for shape in shape_list:
            print("@@@@@@@@@@@@@@@@@@@@",type(shape))
            print(shape.values())
    
    print("\nExtracted Headers:")
    for header in headers:
        print(header)
